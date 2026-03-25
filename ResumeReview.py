# from PySimpleGUI import * (RIP)
from FreeSimpleGUI import *
from docx import Document
from pptx import Presentation
from fitz import open as PDF
import win32com.client as Windows
from pyperclip import *
import openpyxl as xl
from openpyxl.styles import Font
from openpyxl.worksheet.dimensions import ColumnDimension
import os
import re
import inspect
import sys
import datetime as dt
"""
class Format:
    Purple = '\033[95m'
    Cyan = '\033[96m'
    DarkCyan = '\033[36m'
    Blue = '\033[94m'
    Green = '\033[92m'
    Yellow = '\033[93m'
    Red = '\033[91m'
    Bold = '\033[1m'
    Underline = '\033[4m'
    End = '\033[0m'
"""
# Prepping GUI
#theme("DarkGrey5")
theme("DarkTeal12")
folder_selector_column = [
    [
        Text("Select File", font=("Arial", 42))
    ],
    [
        In(size=(35, 1), enable_events=True, key="FolderInput"),
        FolderBrowse(),
        Button("Select", enable_events=True, key="SelectButton")
    ],
    [
        Text("", font=("Arial", 11), enable_events=True, key="MessageBox")
    ],
    [
        Listbox(values=[], enable_events=True, size=(50, 20), key="FilePreview")
    ],
    [
        Button("Continue", enable_events=True, key="FileConfirmation", disabled=True)
    ]
]
review_column = [
    [
        Text("Review", font=("Arial", 42))
    ],
    [
        Multiline(autoscroll=False, enable_events=True, reroute_cprint=True, disabled=True, size=(50, 20), key="ReviewBox")
    ],
    [
        Button("Copy", enable_events=True, key="CopyButton", disabled=True)
    ],
    [
        Text("", font=("Arial", 12), enable_events=True, key="SavedBox")
    ],
    [
        HSeparator()
    ],
    [
        Button("Copy source code - Python3", enable_events=True, key="CopySourceButton")
        #Column([#[Text("Source Code", font=("Arial bold", 15))],
                #[Button("Copy source code - Python3", enable_events=True, key="CopySourceButton")]#,
                #[Button("Copy source code - Other platforms", enable_events=True, key="CopyOtherSourceButton")]
                #]),
        #VSeperator(),
        #Column([#[Text("Note", font=("Arial bold", 15), text_color='yellow')],
                #[Text("Certain features like:\n.doc, .ppt  and .rtf checking are\nonly supported on Windows.", font=("Arial", 8), text_color='yellow')]
                #])
    ]
]
layout = [
    [
        Column(folder_selector_column),
        VSeperator(),
        Column(review_column)
    ]
]
gui = Window("FileKeywordFinder", layout)


def log_keywords():
    global keywordsLen
    global keywords
    global currPara
    global reportSheet
    global y
    r = 0
    for e in range(keywordsLen):
        #occurrences = 0
        if any(keywords[r].endswith(char) or keywords[r].startswith(char) for char in ['#', '*', '$', '@', '<', '>', ':', ';', '/']):
            for k in range(len(re.findall(keywords[r], currPara))):
                #occurrences = occurrences + 1
                reportSheet.cell(row=y + 2, column=r + 3).value = reportSheet.cell(row=y + 2, column=r + 3).value + 1
                #print(occurrences)
                #currPara = re.sub(r"\b" + keywords[r] + r"\b", Format.Bold + Format.Underline + keywords[r] + Format.End, currPara)
        else:
            for k in range(len(re.findall(r"\b" + keywords[r] + r"\b", currPara))):
                #occurrences = occurrences + 1
                reportSheet.cell(row=y + 2, column=r + 3).value = reportSheet.cell(row=y + 2, column=r + 3).value + 1
                #print(occurrences)
                #currPara = re.sub(r"\b" + keywords[r] + r"\b", Format.Bold + Format.Underline + keywords[r] + Format.End, currPara)
        r = r + 1


def format_keywords():
    global keywordsLen
    global keywords
    global currPara
    r = 0
    for e in range(keywordsLen):
        #occurrences = 0
        if any(keywords[r].endswith(char) or keywords[r].startswith(char) for char in ['#', '*', '$', '@', '<', '>', ':', ';', '/']):
            if re.search(keywords[r], currPara) != None:
                #occurrences = occurrences + 1
                #reportSheet.cell(row=y + 2, column=r + 2).value = reportSheet.cell(row=y + 2, column=r + 2).value + 1
                #print(occurrences)
                #currPara = re.sub(r"\b" + keywords[r] + r"\b", Format.Bold + Format.Underline + keywords[r] + Format.End, currPara)
                currPara = re.sub(keywords[r], "[" + keywords[r] + "]", currPara)
        else:
            if re.search(r"\b" + keywords[r] + r"\b", currPara) != None:
                #occurrences = occurrences + 1
                #reportSheet.cell(row=y + 2, column=r + 2).value = reportSheet.cell(row=y + 2, column=r + 2).value + 1
                #print(occurrences)
                #currPara = re.sub(r"\b" + keywords[r] + r"\b", Format.Bold + Format.Underline + keywords[r] + Format.End, currPara)
                currPara = re.sub(r"\b" + keywords[r] + r"\b", "[" + keywords[r] + "]", currPara)
        r = r + 1
    return currPara


def total_hits():
    global reportSheet
    global keywordsLen
    i = 0
    for ff in range(finalFilesListLen):
        j = 2
        totalHits = 0
        for ll in range(keywordsLen):
            totalHits += int(reportSheet.cell(row=i+2, column=j+1).value)
            j = j + 1
        reportSheet.cell(row=i+2, column=2).value = totalHits
        i = i + 1


while True:
    event, values = gui.read()
    if event=="Exit" or event==WIN_CLOSED:
        gui.close()
        break
    elif event=="CopySourceButton":
        code = inspect.getsource(sys.modules[__name__])
        copy(code)
    elif event=="CopyButton":
        copy(values["ReviewBox"])
    elif event=="SelectButton":
        gui["ReviewBox"].update("")
        gui["SavedBox"].update("")
        gui["FileConfirmation"].update(disabled=False)
        gui["CopyButton"].update(disabled=True)
        folderPath = values["FolderInput"]
        folderPath = re.sub('/', '\\\\', folderPath)
        #print(folderPath)
        try:
            filesList = os.listdir(folderPath)
        except FileNotFoundError:
            filesList = []
            gui["MessageBox"].update("Invalid path!", text_color='red')
            gui["FileConfirmation"].update(disabled=True)
        #print(filesList)
        acceptedFileFormats = ["docx", "doc", "pptx", "ppt", "pdf", "rtf"]
        keywords = []
        try:
            keywords = open(folderPath + "\\" + "keywords.txt", "r").read().lower().split(", ")
            #print("LOWERSPLIT", keywords)
        except FileNotFoundError:
            popup("Error", "Couldn't find a valid keywords file.", "Make sure the file name is keywords.txt and that it is stored in the same directory.")
            gui["FileConfirmation"].update(disabled=True)
            #exit()
        keywordsLen = len(keywords)
        #print(keywords)
        initFilesListLen = len(filesList)
        filesListIndexCount = 0
        for z in range(initFilesListLen):
            currFilesListLen = len(filesList)
            fileExtentionFinder = filesList[filesListIndexCount].split(".")
            if fileExtentionFinder[-1] not in acceptedFileFormats:
                #print("non docx file found at index", filesListIndexCount)
                filesList[filesListIndexCount] = " "
                filesListIndexCountCopy = filesListIndexCount
                #print("copy", filesListIndexCountCopy)
                forloopIterationCount = currFilesListLen - (filesListIndexCountCopy + 1)
                #print("forloop iterations", forloopIterationCount)
                for i in range(forloopIterationCount):
                    filesList[filesListIndexCountCopy] = filesList[filesListIndexCountCopy + 1]
                    filesListIndexCountCopy = filesListIndexCountCopy + 1
                del filesList[-1]
                # filesListIndexCount = filesListIndexCount + 1
                #print(filesList)
            elif fileExtentionFinder[-1] in acceptedFileFormats:
                filesListIndexCount = filesListIndexCount + 1
            else:
                popup("Error", "An unknown error occurred.")
        finalFilesListLen = len(filesList)
        gui["FilePreview"].update(filesList)
        if finalFilesListLen>0:
            gui["MessageBox"].update("Found " + str(finalFilesListLen) + " compatible files", text_color='lightgreen')
    elif event=="FileConfirmation":
        gui["CopyButton"].update(disabled=False)
        folderPath = values["FolderInput"]
        folderPath = re.sub('/', '\\\\', folderPath)
        #print(folderPath)
        filesList = os.listdir(folderPath)
        #print(filesList)
        acceptedFileFormats = ["docx", "doc", "pptx", "ppt", "pdf", "rtf"]
        keywords = []
        try:
            keywords = open(folderPath + "\\" + "keywords.txt", "r").read().lower().split(", ")
        except FileNotFoundError:
            popup("Error", "Couldn't find a valid keywords file.", "Make sure the file name is keywords.txt and that it is stored in the same directory.")
            gui["FileConfirmation"].update(disabled=True)
            #exit()
        keywordsLen = len(keywords)
        #print(keywords)
        initFilesListLen = len(filesList)
        filesListIndexCount = 0
        for z in range(initFilesListLen):
            fileExtentionFinder = filesList[filesListIndexCount].split(".")
            if fileExtentionFinder[-1] not in acceptedFileFormats:
                del filesList[filesListIndexCount]
                # filesListIndexCount = filesListIndexCount + 1
                #print(filesList)
            elif fileExtentionFinder[-1] in acceptedFileFormats:
                filesListIndexCount = filesListIndexCount + 1
            else:
                popup("Error", "An unknown error occurred.")
        finalFilesListLen = len(filesList)
        #gui["FilePreview"].update(filesList)
        # continueConfirmation = (str(input("Are these the correct files? y or n: ")))
        count = 0
        y = 0
        reportBook = xl.Workbook()
        reportSheet = reportBook.active
        reportSheet.title = "Report"
        # keywordsXl = [""] + keywords
        #print("EXCEL", keywords)
        reportSheet.append([""] + ["Hits"] + keywords)
        for k in range(finalFilesListLen):
            for b in range(keywordsLen):
                reportSheet.cell(row=k + 2, column=b + 3).value = int(0)
        f = 1
        for j in range(keywordsLen + 2):
            reportSheet.cell(row=1, column=f).font = Font(bold=True)
            f = f + 1
        s = 2
        for h in range(finalFilesListLen):
            reportSheet['a' + str(s)] = filesList[h]
            reportSheet['a' + str(s)].font = Font(bold=True)
            s = s + 1
        try:
            while count < 1:
                if filesList[y].split(".")[-1] == "docx":
                    docxDoc = Document(folderPath + "\\" + filesList[y])
                    # print(folderPath + "\\" + filesList[0])
                    # print(len(docxDoc.paragraphs))
                    # print(y)
                    x = len(docxDoc.paragraphs)
                    i = 0
                    for x in range(x):
                        currPara = docxDoc.paragraphs[i].text
                        # print(docxDoc.paragraphs[i].text)
                        currPara = currPara.lower()
                        i = i + 1
                        log_keywords()
                        cprint(format_keywords())
                elif filesList[y].split(".")[-1] == "doc":
                    word = Windows.Dispatch("Word.Application")
                    word.Visible = True
                    docxDoc = word.Documents.Open(FileName=folderPath + "\\" + filesList[y], Encoding='gbk')
                    # docxDoc = Document(folderPath + "\\" + filesList[y])
                    # print(folderPath + "\\" + filesList[0])
                    # print(len(docxDoc.paragraphs))
                    # print(y)
                    x = len(docxDoc.paragraphs)
                    i = 0
                    for x in range(x):
                        currPara = docxDoc.paragraphs[i].Range.Text
                        # print(docxDoc.paragraphs[i].text)
                        currPara = currPara.lower()
                        i = i + 1
                        log_keywords()
                        cprint(format_keywords())
                    docxDoc.Close()
                    #word.Quit()
                elif filesList[y].split(".")[-1] == "rtf":
                    word = Windows.Dispatch("Word.Application")
                    word.Visible = True
                    docxDoc = word.Documents.Open(FileName=folderPath + "\\" + filesList[y], Encoding='gbk')
                    # docxDoc = Document(folderPath + "\\" + filesList[y])
                    # print(folderPath + "\\" + filesList[0])
                    # print(len(docxDoc.paragraphs))
                    # print(y)
                    x = len(docxDoc.paragraphs)
                    i = 0
                    for x in range(x):
                        currPara = docxDoc.paragraphs[i].Range.Text
                        # print(docxDoc.paragraphs[i].text)
                        currPara = currPara.lower()
                        i = i + 1
                        log_keywords()
                        cprint(format_keywords())
                    docxDoc.Close()
                    #word.Quit()
                elif filesList[y].split(".")[-1] == "pptx":
                    docxDoc = Presentation(folderPath + "\\" + filesList[y])
                    # print(folderPath + "\\" + filesList[0])
                    # print(len(docxDoc.paragraphs))
                    # print(y)
                    i = 0
                    for slide in docxDoc.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                currPara = shape.text
                                # print(docxDoc.paragraphs[i].text)
                                currPara = currPara.lower()
                                i = i + 1
                                log_keywords()
                                cprint(format_keywords())
                elif filesList[y].split(".")[-1] == "ppt":
                    powerpoint = Windows.Dispatch("Powerpoint.Application")
                    powerpoint.Visible = True
                    docxDoc = powerpoint.Presentations.Open(FileName=folderPath + "\\" + filesList[y])
                    # docxDoc = Document(folderPath + "\\" + filesList[y])
                    # print(folderPath + "\\" + filesList[0])
                    # print(len(docxDoc.paragraphs))
                    # print(y)
                    # x = len(docxDoc.paragraphs)
                    # i=0
                    for i in range(1, docxDoc.Slides.Count + 1):
                        for x in range(1, docxDoc.Slides(i).Shapes.Count + 1):
                            if docxDoc.Slides(i).Shapes(x).Hastextframe:
                                paragraphsLen = len(docxDoc.Slides(i).Shapes(x).TextFrame.TextRange.Paragraphs())
                                l = 0
                                for q in range(paragraphsLen):
                                    if l < 1:
                                        l = l + 1
                                        # print("ifloop")
                                        continue
                                    else:
                                        currPara = docxDoc.Slides(i).Shapes(x).TextFrame.TextRange.Paragraphs(
                                            q).text
                                        # print(q)
                                        # print(docxDoc.paragraphs[i].text)
                                        currPara = currPara.lower()
                                        log_keywords()
                                        cprint(format_keywords())
                                currPara = docxDoc.Slides(i).Shapes(x).TextFrame.TextRange.Paragraphs(
                                    paragraphsLen).text
                                # print(docxDoc.paragraphs[i].text)
                                currPara = currPara.lower()
                                log_keywords()
                                cprint(format_keywords())
                    docxDoc.Close()
                    #powerpoint.Quit()
                elif filesList[y].split(".")[-1] == "pdf":
                    docxDoc = PDF(folderPath + "\\" + filesList[y])
                    # print(folderPath + "\\" + filesList[0])
                    # print(len(docxDoc.paragraphs))
                    # print(y)
                    i = 0
                    for page in docxDoc:
                        currPara = page.getText()
                        # print(docxDoc.paragraphs[i].text)
                        currPara = currPara.lower()
                        i = i + 1
                        log_keywords()
                        cprint(format_keywords())
                y = y + 1
                # print(y)
                if y == finalFilesListLen:
                    break
        except KeyboardInterrupt:
            #print('error')
            popup("Error", "Error while reading files")#, "Known glitch: If a .doc file and a .rtf file are next to each other then Word gives an error.")
        total_hits()
        reportSheet.column_dimensions['a'] = ColumnDimension(reportSheet, auto_size=True)
        try:
            reportBook.save(filename=folderPath + "\\" + dt.datetime.today().strftime("Report_%d-%m-%Y_%H.%M.%S.xlsx"))
            gui["SavedBox"].update("The excel file was saved in the directory", text_color='lightgreen')
        except:
            gui["SavedBox"].update("Error saving the excel file, please try again", text_color='red')
        gui["FolderInput"].update("")
        gui["FileConfirmation"].update(disabled=True)
        gui.Element("ReviewBox").set_vscroll_position(0)
        #gui["ReviewBox"].update(reviewContent)